#!/usr/bin/env python3
"""Build the MAC capabilities deck for commit 8b424c20 as a 16:9 .pptx.

This deck is a point-in-time artifact. It is pinned to the commit named in its
directory and is NOT regenerated as the code moves — build a new timestamped
directory instead, so an old deck keeps meaning what it meant when it was shown.

Requires python-pptx, which is deliberately not a repository dependency:

    python3 -m venv /tmp/deckvenv
    /tmp/deckvenv/bin/pip install python-pptx
    /tmp/deckvenv/bin/python build_deck.py

Diagram PNGs are rendered from the SVGs beside them; see README.md.
"""

from __future__ import annotations

from pathlib import Path
from struct import unpack

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
IMAGES = HERE / "images"
COMMIT = "8b424c20"
CAPTURED = "2026-08-20T01:12:24Z"
OUT = HERE / f"mac-capabilities-{COMMIT}.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x1B, 0x25, 0x30)
SLATE = RGBColor(0x2F, 0x3A, 0x45)
GREY = RGBColor(0x5A, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x98, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0x9E)
AMBER = RGBColor(0xB5, 0x65, 0x1D)
SAND = RGBColor(0xF2, 0xC4, 0x8A)
PALE = RGBColor(0xDC, 0xE5, 0xEC)
GREEN = RGBColor(0x26, 0x55, 0x36)

FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def bg(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    return tf


def line(tf, text, size, color, bold=False, space_before=0, space_after=6,
         italic=False, first=False, mono=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Menlo" if mono else FONT
    return p


def add_image_slide(name: str, note: str):
    """Full-bleed diagram slide: each diagram carries its own title."""
    slide = prs.slides.add_slide(BLANK)
    path = IMAGES / name
    with open(path, "rb") as fh:
        head = fh.read(33)
    iw, ih = unpack(">II", head[16:24])
    aspect = iw / ih

    h, max_w = Inches(7.1), Inches(12.9)
    w = Emu(int(h * aspect))
    if w > max_w:
        w = max_w
        h = Emu(int(w / aspect))
    slide.shapes.add_picture(
        str(path), Emu(int((SLIDE_W - w) / 2)), Emu(int((SLIDE_H - h) / 2)),
        width=w, height=h,
    )
    notes(slide, note)
    return slide


def section(kicker: str, title: str, subtitle: str, note: str) -> None:
    slide = prs.slides.add_slide(BLANK)
    bg(slide, SLATE)
    tf = textbox(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.4))
    line(tf, kicker, 16, SAND, bold=True, first=True, space_after=10)
    line(tf, title, 36, WHITE, bold=True, space_after=12)
    line(tf, subtitle, 18, PALE)
    notes(slide, note)


# ------------------------------------------------------------------ slides

def title_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    bg(slide, SLATE)
    tf = textbox(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.6))
    line(tf, "MAC", 54, WHITE, bold=True, first=True, space_after=6)
    line(tf, "What the control plane can do today", 30, PALE, space_after=18)
    line(tf, "A multi-agent coordinator control plane, audited from its own source and documentation.",
         17, MUTED)

    tf2 = textbox(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.4))
    line(tf2, f"commit {COMMIT}   ·   captured {CAPTURED}", 16, SAND, bold=True, first=True,
         space_after=4)
    line(tf2, "Every claim in this deck is traced to a file or commit in AUDIT.md. Figures are dated because they age.",
         13.5, MUTED)

    notes(slide, f"""
This deck describes MAC at commit {COMMIT} and nothing later. It was built by auditing the
source and docs directly rather than from memory, and AUDIT.md traces every number back to the
file or commit it came from.

Open by saying what MAC is NOT: it is not a chat agent, not an IDE, and not a model. It is the
durable operational record underneath a fleet of coding agents.
""")


def what_it_is() -> None:
    slide = prs.slides.add_slide(BLANK)
    tf = textbox(slide, Inches(0.85), Inches(0.55), Inches(11.6), Inches(1.0))
    line(tf, "What MAC is, and what it deliberately is not", 34, INK, bold=True, first=True,
         space_after=4)
    line(tf, "The boundary is the design. Everything MAC owns is something you would want to still be true after a crash.",
         15, GREY)

    tf2 = textbox(slide, Inches(0.85), Inches(1.85), Inches(5.6), Inches(5.0))
    line(tf2, "It owns durable operational truth", 20, BLUE, bold=True, first=True, space_after=10)
    for t in [
        "Tasks, leases, routing, reviews, evidence.",
        "Secrets as handles, runtime manifests, rollout state.",
        "Machine and agent registry, capabilities and health.",
        "Audit trails, and the publication record.",
    ]:
        line(tf2, "•  " + t, 15, SLATE, space_after=8)
    line(tf2, "PostgreSQL is the authority — there is no offline replica and no ticket-file sync. "
              "A .tickets/<id>.md file is an optional local convenience view and is explicitly not authoritative.",
         14, GREY, space_before=10)

    tf3 = textbox(slide, Inches(7.0), Inches(1.85), Inches(5.5), Inches(5.0))
    line(tf3, "It does not own the conversation", 20, AMBER, bold=True, first=True, space_after=10)
    for t in [
        "Personality, chat, and messaging gateways belong to the human-facing runtime.",
        "MAC sits underneath it and takes durable work from it.",
        "The command-and-control dashboard was retired; /ui is read-only observability.",
        "The console speaks on the bus like any other participant rather than mutating the control plane.",
    ]:
        line(tf3, "•  " + t, 15, SLATE, space_after=8)
    line(tf3, "That boundary is why an agent going off the rails is a bus conversation, "
              "not a privileged API call.", 14, GREY, space_before=10)

    notes(slide, """
The line to land: MAC owns what must survive a crash; the runtime above it owns what must feel
like a conversation.

The retirement of the command-and-control dashboard is worth calling out — it is the clearest
signal of the direction of travel. The console can no longer mutate the control plane at all.
""")


def scale_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    tf = textbox(slide, Inches(0.85), Inches(0.55), Inches(11.6), Inches(1.0))
    line(tf, "Scale, at this commit", 34, INK, bold=True, first=True, space_after=4)
    line(tf, f"Counted directly in the tree at {COMMIT}, not estimated.", 15, GREY)

    rows = [
        ("195,764", "lines of Python under src/", "201 modules in src/mac"),
        ("470", "test files under tests/", "plus fault-replay and contract suites"),
        ("408", "HTTP routes", "generated from the live OpenAPI schema"),
        ("123", "CLI verbs", "across project, task, agent and admin"),
        ("18", "book chapters", "every shell example executed by make docs-check"),
        ("18", "architecture decision records", "0016–0018 are Proposed, not shipped"),
        ("5", "coding-agent routes", "claude, codex, cursor, opencode, pi"),
    ]
    tf2 = textbox(slide, Inches(0.85), Inches(1.8), Inches(11.6), Inches(5.2))
    first = True
    for num, label, note_ in rows:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        p.space_after = Pt(11)
        r1 = p.add_run(); r1.text = f"{num:>10}   "
        r1.font.size = Pt(23); r1.font.bold = True; r1.font.color.rgb = INK; r1.font.name = FONT
        r2 = p.add_run(); r2.text = label
        r2.font.size = Pt(17); r2.font.color.rgb = SLATE; r2.font.name = FONT
        r3 = p.add_run(); r3.text = f"   — {note_}"
        r3.font.size = Pt(14); r3.font.color.rgb = GREY; r3.font.name = FONT
        first = False

    notes(slide, """
Do not read this table aloud. Use it to make one point: this is not a prototype, and the
documentation is generated from the running system rather than written alongside it.

The generated CLI and OpenAPI references are checked in CI, so a drifted doc fails the build.
""")


def honesty_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    tf = textbox(slide, Inches(0.85), Inches(0.55), Inches(11.6), Inches(1.0))
    line(tf, "Proposed, deferred, or stale — stated up front", 34, INK, bold=True, first=True,
         space_after=4)
    line(tf, "A capabilities deck that only lists capabilities is marketing. These are the places the code and the claims do not yet meet.",
         15, GREY)

    items = [
        ("ADR 0016, 0017 and 0018 are all Proposed.",
         "Agent-initiated review, router-side token metering, and the dependency graph are designs with measured motivation — none of them has shipped. The mandatory hub-mediated workflow still runs today."),
        ("ADR 0012 is Accepted with implementation deferred.",
         "The native-steward / containerized-execution split is a decision awaiting fleet measurement, and ADR 0015 has already narrowed its containerized half to Linux nodes only."),
        ("The README still describes a runtime that was deleted.",
         "It documents src/mac/_hermes as a vendored Hermes snapshot, links a snapshot contract that no longer resolves, and claims coverage and lint exclusions for that path. The tree went in 3ebde2dd. The mac-hermes adapter it also names is real and unaffected."),
        ("Token attribution is not currently trustworthy.",
         "29.5% of routes over the last seven days recorded no input token count, none reported cache hits, and cost is priced at read time rather than stored. ADR 0017 exists because of this."),
        ("A third of blocked work is dead and nothing surfaces it.",
         "165 of 355 blocked tasks wait on dependencies that can never complete. Finding that took a hand-written query."),
    ]
    tf2 = textbox(slide, Inches(0.85), Inches(1.8), Inches(11.6), Inches(5.3))
    first = True
    for head, body in items:
        line(tf2, head, 17.5, INK, bold=True, space_before=0 if first else 10, space_after=3,
             first=first)
        line(tf2, body, 14, GREY, space_after=2)
        first = False

    notes(slide, """
This is the slide that earns the rest of the deck. In front of engineers, naming the five
weakest points yourself is what makes the capability claims credible.

The README staleness is a real finding from this audit and is worth fixing separately.
""")


def provenance_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    bg(slide, SLATE)
    tf = textbox(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
    line(tf, "Provenance", 34, WHITE, bold=True, first=True, space_after=4)
    line(tf, "This deck is a pinned artifact, not a living document.", 16, PALE)

    tf2 = textbox(slide, Inches(0.9), Inches(2.0), Inches(5.6), Inches(4.6))
    line(tf2, "How to read it", 19, SAND, bold=True, first=True, space_after=8)
    for t in [
        f"It describes commit {COMMIT} and nothing later.",
        "Figures carry the date they were measured, because they age.",
        "AUDIT.md traces every claim to a file, commit or generated reference.",
        "Diagrams are SVG sources; the PNGs are rendered from them.",
    ]:
        line(tf2, "•  " + t, 14.5, PALE, space_after=8)

    tf3 = textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    line(tf3, "Making the next one", 19, SAND, bold=True, first=True, space_after=8)
    line(tf3, "Create a new sibling directory rather than editing this one:", 14.5, PALE,
         space_after=8)
    line(tf3, "docs/presentation/<UTC timestamp>-<short commit>/", 13.5, WHITE, mono=True,
         space_after=10)
    line(tf3, "The timestamp sorts, the commit pins, and neither collides. An old deck keeps "
              "meaning what it meant when it was shown.", 14.5, PALE, space_after=8)
    line(tf3, "See README.md in this directory for the exact render and build commands.",
         14.5, MUTED)

    notes(slide, """
Close on the convention rather than a summary: decks are cheap to regenerate and are pinned to a
commit, so nobody has to wonder whether a slide is still true — they can check the hash.
""")


# ------------------------------------------------------------------ build

title_slide()
what_it_is()

section("PART ONE", "The model",
        "Four objects, one ledger, and what must be true to move work between states.",
        "Transition: everything MAC does is expressed against four nouns.")

add_image_slide("01-object-model.png", """
The CLI is not a wrapper over the object model — it IS the object model. project, task, agent,
admin, 123 verbs.

Two details worth pointing at: recovery is a first-class set of verbs rather than an incident
runbook, and break-glass is grantable, listable and revocable, so an emergency leaves a record
instead of an ambient permission.

The bottom band matters for integrators: six surfaces, one dispatch seam, and a contract test
that proves the Python client and the hub agree about the route table.
""")

add_image_slide("02-task-lifecycle.png", """
Eleven states, three terminal. The default list view shows non-terminal work, because the useful
question is "what still wants something from somebody".

The four gates are the real content: a lease and fence authorize every mutation, evidence is typed
and bound to one exact attempt, review must come from a different agent AND a different model, and
landing goes through a speculative merge queue whose invariant is "never land an untested tree".

If asked why MAC built its own merge queue: GitHub's is organization-only, verified by API against
the live repo — HTTP 422 on a user-owned repository.
""")

section("PART TWO", "Coordination and execution",
        "How agents hear each other, and what actually runs the work.",
        "Transition: the coordination model changed shape recently, and deliberately.")

add_image_slide("03-coordination.png", """
The headline: AgentBus became a broadcast bus. Point-to-point messages are no longer private,
because an agent told "worker-2 already rebased that branch" cannot verify it without reading
worker-2's stream.

The motivating incidents are real and on the record — a `git add -A` moments from sweeping ~1,200
lines of another agent's work into an unrelated commit, and a `git commit -a` that did.

Note the human is a participant on the bus, not a controller above it, and that stand_down and
abort are deliberately separate verbs. One carve-out: raw terminal streams stay participant-scoped
because they carry credentials.
""")

add_image_slide("04-fleet-execution.png", """
MAC does not pretend the fleet is homogeneous. macOS nodes are host installs under launchd — an
earlier attempt to containerize them was superseded. Linux nodes run a native steward with
containerized execution. Kubernetes and HGX provide elastic capacity.

Five coding-agent routes, and the sandbox derives which are available from the router rather than
hardcoding a list.

The identity band is the one auditors care about: secrets are handles with audit records, output
is redacted, and egress is declared per project and per task.
""")

section("PART THREE", "Measurement",
        "The fleet instruments itself — and the instrumentation is currently indicting the design.",
        "Transition: the most interesting capability is that MAC can prove itself wrong.")

add_image_slide("05-measurement.png", """
THE KEY SLIDE. Left: what MAC measures — review experiments with pre-registered arms and a blind
treatment, a scientific optimizer that only ever emits a policy candidate, a dreaming pipeline
where memory must SHRINK to be promoted, and operational learning that changes routing.

Right: what those measurements actually found, dated. 29.5% of model routes recorded no input
tokens; none reported cache hits. 165 of 355 blocked tasks are permanently dead.

Bottom: three open ADRs, all raised in the last three days, each caused by a measurement rather
than an opinion. That is the point of the slide — the instrumentation is what generated the
backlog.
""")

scale_slide()
honesty_slide()
provenance_slide()

prs.save(OUT)
print(f"wrote {OUT.name} — {len(prs.slides._sldIdLst)} slides")
