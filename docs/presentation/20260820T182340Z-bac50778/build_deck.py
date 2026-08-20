#!/usr/bin/env python3
"""Build the MAC architecture deck for commit bac50778 as a 16:9 .pptx.

Deliberately picture-first: three flow diagrams and three live console captures
carry the argument, and the text slides only hold what a picture cannot say —
counts, and what has not shipped.

The console captures are NOT reproducible from this repository. They come from a
running hub, via the command in README.md, and are gitignored along with every
other PNG here. Rebuilding this deck without a live fleet gives you the diagrams
and blank screenshot slides.

Requires python-pptx, which is deliberately not a repository dependency:

    python3 -m venv /tmp/deckvenv && /tmp/deckvenv/bin/pip install python-pptx
    /tmp/deckvenv/bin/python build_deck.py

Then publish with scripts/publish-deck-to-slides.py. See skills/cut-a-release.
"""

from __future__ import annotations

from pathlib import Path
from struct import unpack

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
IMAGES = HERE / "images"
COMMIT = "bac50778"
CAPTURED = "2026-08-20T18:23:40Z"
OUT = HERE / f"mac-architecture-{COMMIT}.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x1B, 0x25, 0x30)
SLATE = RGBColor(0x2F, 0x3A, 0x45)
GREY = RGBColor(0x5A, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x98, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0x9E)
AMBER = RGBColor(0xB5, 0x65, 0x1D)
SAND = RGBColor(0xF2, 0xC4, 0x8A)
PALE = RGBColor(0xDC, 0xE5, 0xEC)
FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    return tf


def line(tf, text, size, color, bold=False, space_before=0, space_after=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before, p.space_after = Pt(space_before), Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    run.font.name = FONT
    return p


def _aspect(path: Path) -> float:
    with open(path, "rb") as fh:
        head = fh.read(33)
    iw, ih = unpack(">II", head[16:24])
    return iw / ih


def image_slide(name: str, note: str, caption: str | None = None, sub: str | None = None):
    """Full-bleed when the diagram carries its own title; captioned for captures."""
    slide = prs.slides.add_slide(BLANK)
    path = IMAGES / name
    if not path.exists():
        notes(slide, note)
        return slide

    top_pad = Inches(1.15) if caption else Inches(0.2)
    if caption:
        tf = textbox(slide, Inches(0.6), Inches(0.32), Inches(12.2), Inches(0.9))
        line(tf, caption, 27, INK, bold=True, first=True, space_after=2)
        if sub:
            line(tf, sub, 14, GREY)

    avail_h = SLIDE_H - top_pad - Inches(0.25)
    avail_w = Inches(12.6)
    aspect = _aspect(path)
    h, w = avail_h, Emu(int(avail_h * aspect))
    if w > avail_w:
        w, h = avail_w, Emu(int(avail_w / aspect))
    slide.shapes.add_picture(
        str(path), Emu(int((SLIDE_W - w) / 2)), Emu(int(top_pad + (avail_h - h) / 2)),
        width=w, height=h,
    )
    notes(slide, note)
    return slide


# ------------------------------------------------------------------ slides

slide = prs.slides.add_slide(BLANK)
bg(slide, SLATE)
tf = textbox(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
line(tf, "MAC", 54, WHITE, bold=True, first=True, space_after=6)
line(tf, "How the control plane is put together", 30, PALE, space_after=18)
line(tf, "One hub, many workers, and what runs inside the hub.", 17, MUTED)
tf2 = textbox(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2))
line(tf2, f"commit {COMMIT}   ·   captured {CAPTURED}", 16, SAND, bold=True, first=True, space_after=4)
line(tf2, "Diagrams from the source. Screenshots from a hub that was running while this was made.", 13.5, MUTED)
notes(slide, """
This deck is deliberately picture-first: three flow diagrams and three live console captures.
It describes commit bac50778 and nothing later.

Open with the shape, not the philosophy: one hub, many workers, and the hub is smaller than people
expect.
""")

image_slide("01-hub-and-workers.png", """
THE CORE RELATIONSHIP. The hub owns three things — the task ledger, the capability registry, and
capacity — and hands out leases. Workers hold the actual work.

The amber box is the point: the hub does NOT decide what a task needs. That belongs to the agent
that read the diff, because it is the only party that can know.

Note the worker classes are genuinely different: macOS is a host install under launchd, Linux runs
containerized execution under a native steward, Kubernetes/HGX is elastic. MAC does not pretend
they are one server class.

AgentBus runs underneath all of it, and the human is ON it — the console never mutates the control
plane, it speaks like any other participant.

Live numbers bottom-right: 10 agents, 2 busy / 2 idle / 6 offline. Offline is normal — fungible
workers are created on demand.
""")

image_slide("03-life-of-a-task.png", """
ONE TASK, END TO END. Ten steps, five lanes, time left to right.

The thing to say out loud: every hand-off is a durable row, not a message. That is what makes
"why is this stuck?" answerable months later.

The two red dashed edges are the failure paths, and they matter more than the happy path — lease
expiry reclaims work from a dead agent, and a merge-queue eviction discards everything tested
behind the evicted entry because those results were green against a state that will never exist.

Step 6 is changing: ADR 0016 was accepted today, making review agent-initiated rather than
mandatory. The motivating measurement is on the slide — 52 reviews, zero findings.
""")

image_slide("02-inside-the-hub.png", """
THE DEEP DIVE, and the slide to spend time on.

Three daemon threads, each separate for a measured reason. mac-publication got its own thread
because it used to run inside the tick AND on worker heartbeats — heartbeats were measured at
250-315 seconds because the cost landed on the worker.

The middle band is the tick, in order. The order IS the design: retention sits at position 7, ahead
of the review sweep, because it used to sit behind it and was starved — zero retention events in 48
minutes while 235,615 rows sat past the cutoff. That failure is invisible by construction, since
retention is silent when idle.

31 services grouped by purpose. Do not read the lists; use them to show the surface is organised
rather than accreted.
""")

image_slide(
    "04-console-fleet.png",
    """
The read-only observability console, on a live hub.

Ten agents; the "status not believable" tile is the interesting one — it counts agents reporting
idle or busy that have not been heard from in over 15 minutes. A fleet that trusts self-reported
status has no way to notice that.

Agent identity is cropped out deliberately: this repository's docs must read as generic for any
fleet owner, and the roster names would fail that gate.
""",
    caption="The fleet, observed",
    sub="/ui is read-only. The command-and-control dashboard was retired — the console cannot mutate the control plane.",
)

image_slide(
    "05-console-live.png",
    """
The live view: task movement over a six-hour window, in-flight work by state, and a transition
ticker of individual state changes.

Point at the blocked bar — 465 of 644 in-flight tasks. That is the number ADR 0018 exists for: the
console can say how many are blocked but not WHAT blocks them, because the hub does not ship the
dependency edges.

The ticker is the honest bit: you can watch claimed -> running -> needs_review -> reviewing happen
without anyone driving it.
""",
    caption="Movement, not a snapshot",
    sub="Task transitions over six hours on a live hub, with in-flight work by state.",
)

image_slide(
    "06-console-merge-queue.png",
    """
The merge queue, and the most uncomfortable slide in the deck.

Two queues, 4 waiting, 2 landed — and 211 evicted, for a 1% land rate. Every eviction reason is the
same: "projected merge conflicts with the queue base".

That is not a queue working; that is a queue thrashing against a trunk moving faster than entries
can be tested. It is exactly the kind of thing a dashboard is for, and exactly the kind of thing
prose in a README would never have surfaced.

Repository names are visible and that is fine — the identity gate explicitly permits the repo-org
slug, and forbids only operator and fleet identity.
""",
    caption="Where changes actually land — and where they do not",
    sub="Two queues, a speculative window, and 211 evictions at a 1% land rate.",
)

slide = prs.slides.add_slide(BLANK)
tf = textbox(slide, Inches(0.85), Inches(0.55), Inches(11.6), Inches(1.0))
line(tf, "Scale at this commit", 34, INK, bold=True, first=True, space_after=4)
line(tf, f"Counted in the tree at {COMMIT}, not estimated.", 15, GREY)
rows = [
    ("198,038", "lines of Python under src/", "205 modules"),
    ("483", "test files", "plus fault-replay and contract suites"),
    ("408", "HTTP routes", "generated from the live OpenAPI schema"),
    ("124", "CLI verbs", "across project, task, agent and admin"),
    ("31", "hub services", "grouped six ways on the deep-dive slide"),
    ("24", "architecture decision records", "11 of them still Proposed"),
]
tf2 = textbox(slide, Inches(0.85), Inches(1.9), Inches(11.6), Inches(5.0))
first = True
for num, label, note_ in rows:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph()
    p.space_after = Pt(13)
    r1 = p.add_run(); r1.text = f"{num:>9}   "
    r1.font.size, r1.font.bold, r1.font.color.rgb, r1.font.name = Pt(25), True, INK, FONT
    r2 = p.add_run(); r2.text = label
    r2.font.size, r2.font.color.rgb, r2.font.name = Pt(18), SLATE, FONT
    r3 = p.add_run(); r3.text = f"   — {note_}"
    r3.font.size, r3.font.color.rgb, r3.font.name = Pt(14), GREY, FONT
    first = False
notes(slide, """
Do not read the table. One point: this is not a prototype, and the CLI and HTTP references are
generated from the running system, so they cannot drift without failing CI.
""")

slide = prs.slides.add_slide(BLANK)
tf = textbox(slide, Inches(0.85), Inches(0.55), Inches(11.6), Inches(1.0))
line(tf, "What is decided but not yet true", 34, INK, bold=True, first=True, space_after=4)
line(tf, "The easiest claims in this deck to falsify, stated up front.", 15, GREY)
items = [
    ("Eleven of twenty-four ADRs are still Proposed.",
     "Including router-side token metering (0017) and the dependency graph (0018). A proposal is a decision the code has not caught up with."),
    ("ADR 0016 was accepted the day this deck was built.",
     "Agent-initiated review is decided, not deployed: the rollout is one project-level flag, and the mandatory path still runs everywhere else."),
    ("A 1% merge-queue land rate is on a slide, not in a runbook.",
     "211 evictions, all for projected conflicts with the queue base. Nobody has fixed it; the console is simply honest enough to show it."),
    ("A third of blocked work may be permanently dead.",
     "ADR 0018 found 165 of 355 blocked tasks waiting on dependencies that can never complete, and it took a hand-written query to find."),
    ("The screenshots cannot be regenerated from this repository.",
     "They come from a live hub. The diagrams are reproducible; the captures are evidence with a date on them."),
]
tf2 = textbox(slide, Inches(0.85), Inches(1.85), Inches(11.6), Inches(5.3))
first = True
for head, body in items:
    line(tf2, head, 17.5, INK, bold=True, space_before=0 if first else 11, space_after=3, first=first)
    line(tf2, body, 14, GREY, space_after=2)
    first = False
notes(slide, """
Keep this slide. In front of engineers it is what earns the rest of the deck — and every item is
something the system itself surfaced rather than something a reviewer had to dig out.
""")

slide = prs.slides.add_slide(BLANK)
bg(slide, SLATE)
tf = textbox(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
line(tf, "Provenance", 34, WHITE, bold=True, first=True, space_after=4)
line(tf, "A pinned artifact, not a living document.", 16, PALE)
tf2 = textbox(slide, Inches(0.9), Inches(2.0), Inches(5.6), Inches(4.6))
line(tf2, "How to read it", 19, SAND, bold=True, first=True, space_after=8)
for t in [
    f"It describes commit {COMMIT} and nothing later.",
    "Figures carry the date they were measured.",
    "AUDIT.md traces every claim to a file or commit.",
    "Diagrams are SVG; PNGs render from them.",
]:
    line(tf2, "•  " + t, 14.5, PALE, space_after=8)
tf3 = textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
line(tf3, "Making the next one", 19, SAND, bold=True, first=True, space_after=8)
line(tf3, "Follow skills/cut-a-release/SKILL.md. Build a new sibling directory rather than editing this one:",
     14.5, PALE, space_after=8)
line(tf3, "docs/presentation/<UTC timestamp>-<short commit>/", 13.5, WHITE, space_after=10)
line(tf3, "The timestamp sorts, the commit pins, and neither collides. An old deck keeps meaning what it meant when it was shown.",
     14.5, PALE)
notes(slide, """
Close on the convention: decks are cheap, pinned to a commit, and nobody has to wonder whether a
slide is still true — they can check the hash.
""")

prs.save(OUT)
print(f"wrote {OUT.name} — {len(prs.slides._sldIdLst)} slides")
