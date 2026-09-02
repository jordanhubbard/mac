#!/usr/bin/env python3
"""Build the six-slide MAC v1.3.5 release-candidate deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
COMMIT = "a168e9d0"
CAPTURED = "2026-09-02T13:13:14Z"
OUT = HERE / f"mac-capabilities-{COMMIT}.pptx"

INK = RGBColor(0x1B, 0x25, 0x30)
SLATE = RGBColor(0x2F, 0x3A, 0x45)
GREY = RGBColor(0x5A, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x98, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0x9E)
AMBER = RGBColor(0xB5, 0x65, 0x1D)
SAND = RGBColor(0xF2, 0xC4, 0x8A)
PALE = RGBColor(0xDC, 0xE5, 0xEC)
FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def background(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.paragraphs[0].alignment = align
    return frame


def line(frame, text, size, color, *, bold=False, first=False, after=8):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(after)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def title_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.8))
    line(frame, "MAC", 54, WHITE, bold=True, first=True, after=8)
    line(frame, "v1.3.5 release candidate", 32, PALE, after=16)
    line(frame, "A release procedure that makes the evidence, documentation, and rollout visible.", 18, MUTED)
    footer = text_box(slide, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.8))
    line(footer, f"commit {COMMIT}   ·   captured {CAPTURED}", 16, SAND, bold=True, first=True)
    notes(slide, "This deck is pinned to the release candidate. AUDIT.md traces each visible claim to the source tree or generated reference.")


def change_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(frame, "This release makes long-running change safer to observe and finish", 33, INK, bold=True, first=True)
    line(frame, "The goal is operational confidence, not a faster-looking dashboard.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.05), Inches(11.3), Inches(4.65))
    for index, (head, detail) in enumerate([
        ("Artifact publication", "A tag now builds and validates a versioned wheel before publishing the release."),
        ("Deploy resilience", "A slow but healthy gateway is recorded as degraded rather than turning a node rollout into a false failure."),
        ("Fleet visibility", "AgentBus traffic, roll-call, and fleet news give operators durable read-side evidence."),
        ("Healthy test budget", "The contract suite has up to one hour when it is making progress, so verification is not mistaken for a hang."),
    ]):
        line(body, head, 21, BLUE if index % 2 == 0 else AMBER, bold=True, first=index == 0, after=2)
        line(body, detail, 17, SLATE, after=12)
    notes(slide, "Each point is a shipped change in the audited range, not a planned capability.")


def release_flow_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(frame, "A release is now a gated sequence with a clear stop point", 33, INK, bold=True, first=True)
    line(frame, "The command cannot tag first and explain later.", 16, GREY)
    body = text_box(slide, Inches(1.0), Inches(2.1), Inches(11.1), Inches(4.5), align=PP_ALIGN.CENTER)
    steps = ["Audited docs and deck", "Lint, tests, and docs gate", "Release PR", "Tag and artifacts", "Optional fleet rollout"]
    for idx, step in enumerate(steps):
        line(body, f"{idx + 1}.  {step}", 23, BLUE if idx < 3 else AMBER, bold=True, first=idx == 0, after=14)
    notes(slide, "The release target enforces this ordering. Deployment remains optional and follows artifact publication.")


def evidence_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
    line(frame, "The candidate has a verifiable operating surface", 34, WHITE, bold=True, first=True)
    body = text_box(slide, Inches(0.95), Inches(2.0), Inches(11.2), Inches(4.3))
    for idx, item in enumerate([
        "433 HTTP routes in the generated OpenAPI reference",
        "Six top-level command groups in the generated CLI reference",
        "18 executable documentation chapters",
        "Eight commits since v1.3.4",
        "Expanded sanity scope passed before the release workflow was merged",
    ]):
        line(body, item, 22, PALE if idx < 3 else SAND, bold=idx < 3, first=idx == 0, after=12)
    notes(slide, "Counts are dated observations at this commit. AUDIT.md gives the collection source for each one.")


def boundary_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(frame, "Release automation does not erase operational judgement", 33, INK, bold=True, first=True)
    line(frame, "Automation proves prerequisites; it does not conceal a failed rollout.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.1), Inches(11.3), Inches(4.5))
    line(body, "If the final version bump is not green, the process stops before a tag.", 22, BLUE, bold=True, first=True, after=15)
    line(body, "If deployment fails, the newest generation is retained with diagnostic evidence and a safety hold for a forward repair.", 22, AMBER, bold=True, after=15)
    line(body, "The release record, the deck, and the fleet outcome stay separately inspectable.", 22, SLATE, bold=True)
    notes(slide, "This is a fix-forward policy. A rollback requires explicit break-glass authority.")


def closing_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.7), align=PP_ALIGN.CENTER)
    line(frame, "v1.3.5 is ready to be proven, published, and rolled out", 34, WHITE, bold=True, first=True, after=14)
    line(frame, "The release target keeps the proof with the release instead of leaving it in a terminal transcript.", 20, PALE)
    notes(slide, "Close with the practical consequence: a release is now a reproducible, reviewable procedure rather than a sequence remembered by one operator.")


for build in (title_slide, change_slide, release_flow_slide, evidence_slide, boundary_slide, closing_slide):
    build()

prs.save(OUT)
print(OUT)
