#!/usr/bin/env python3
"""Build the six-slide MAC v1.3.5 release deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
COMMIT = "c7a3fee1"
CAPTURED = "2026-09-04T21:25:15Z"
OUT = HERE / f"mac-capabilities-{COMMIT}.pptx"

INK = RGBColor(0x1B, 0x25, 0x30)
SLATE = RGBColor(0x2F, 0x3A, 0x45)
GREY = RGBColor(0x5A, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x98, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0x9E)
AMBER = RGBColor(0xB5, 0x65, 0x1D)
SAND = RGBColor(0xF2, 0xC4, 0x8A)
PALE = RGBColor(0xDC, 0xE5, 0xEC)
FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def background(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.paragraphs[0].alignment = align
    return frame


def line(frame, text, size, color, *, bold=False, first=False, after=8):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(after)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def title_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.8))
    line(frame, "MAC", 54, WHITE, bold=True, first=True, after=8)
    line(frame, "v1.3.5", 32, PALE, after=16)
    line(
        frame,
        "Root-causing and permanently fixing onboarding, dispatch, and attestation bugs the fleet had been living with.",
        18,
        MUTED,
    )
    footer = text_box(slide, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.8))
    line(footer, f"commit {COMMIT}   ·   captured {CAPTURED}", 16, SAND, bold=True, first=True)
    notes(
        slide,
        "This deck is pinned to the release commit. AUDIT.md traces each visible claim to the source tree or generated reference.",
    )


def onboarding_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(
        frame,
        "OpenShell/OpenClaw onboarding no longer fails on healthy nodes",
        33,
        INK,
        bold=True,
        first=True,
    )
    line(frame, "Four separate false-failure and mismatch bugs, fixed at the root.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.05), Inches(11.3), Inches(4.65))
    for index, (head, detail) in enumerate(
        [
            (
                "Reviewed-CLI preflight",
                "Full identity is computed whenever a canonical CLI binary already exists, not only on first install.",
            ),
            (
                "Degraded-gateway tolerance",
                "A missing service-advertisement file under a tolerated degraded gateway no longer crashes node install.",
            ),
            (
                "Cold-pull gateway race",
                "Bootstrap polls for local gateway readiness for up to 120s instead of a fixed 3-second sleep.",
            ),
            (
                "Static sandbox binary",
                "openshell-sandbox is verified statically linked before install, eliminating host/container glibc mismatches.",
            ),
        ]
    ):
        line(
            body, head, 21, BLUE if index % 2 == 0 else AMBER, bold=True, first=index == 0, after=2
        )
        line(body, detail, 17, SLATE, after=12)
    notes(slide, "Each point is a shipped fix in the audited range, traced in AUDIT.md.")


def dispatch_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(
        frame,
        "Targeted dispatch and fleet recovery are now reliable",
        33,
        INK,
        bold=True,
        first=True,
    )
    line(frame, "An agent explicitly targeted for a task now reliably claims it.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.05), Inches(11.3), Inches(4.65))
    for index, (head, detail) in enumerate(
        [
            (
                "Explicit target claim",
                "A task with target_agent_id is claimed directly instead of depending only on the global dispatch round.",
            ),
            (
                "Attestation reconciliation",
                "retain_forward node recovery reconciles the bound worker's attestation authority against the correct hub.",
            ),
            (
                "Honest sandbox status",
                "agent_status now requires a live sandbox_id before reporting deployed, closing a false-healthy report.",
            ),
            (
                "Ambient-drift immunity",
                "OPENSHELL_GATEWAY_ENDPOINT is pinned into mac.env, so mac-agent no longer depends on ambient gateway selection.",
            ),
        ]
    ):
        line(
            body, head, 21, BLUE if index % 2 == 0 else AMBER, bold=True, first=index == 0, after=2
        )
        line(body, detail, 17, SLATE, after=12)
    notes(slide, "Each point closes a bug that had been worked around operationally, not fixed.")


def evidence_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
    line(frame, "The release has a verifiable operating surface", 34, WHITE, bold=True, first=True)
    body = text_box(slide, Inches(0.95), Inches(2.0), Inches(11.2), Inches(4.3))
    for idx, item in enumerate(
        [
            "433 HTTP routes in the generated OpenAPI reference",
            "Six top-level command groups in the generated CLI reference",
            "38 commits since v1.3.4, all traced in AUDIT.md",
            "11,438 tests passed locally; CI green on the release ancestor",
            "Two pre-existing, unrelated CI flakes identified and judged non-blocking",
        ]
    ):
        line(body, item, 22, PALE if idx < 3 else SAND, bold=idx < 3, first=idx == 0, after=12)
    notes(
        slide,
        "Counts are dated observations at this commit. AUDIT.md gives the collection source for each one.",
    )


def boundary_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(
        frame,
        "These were root-cause fixes, not new capabilities",
        33,
        INK,
        bold=True,
        first=True,
    )
    line(frame, "No CLI flag, HTTP route, or documented contract changed.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.1), Inches(11.3), Inches(4.5))
    line(
        body,
        "Every fix in this range closes a gap between documented/intended behavior and what the fleet actually did.",
        22,
        BLUE,
        bold=True,
        first=True,
        after=15,
    )
    line(
        body,
        "Nothing in README.md or the generated CLI/OpenAPI references needed updating as a result.",
        22,
        AMBER,
        bold=True,
        after=15,
    )
    line(
        body,
        "Two of the fixes (openshell-sandbox linking, targeted dispatch) were authored directly by fleet agents working the reopened tasks.",
        22,
        SLATE,
        bold=True,
    )
    notes(slide, "See AUDIT.md's release-claims table for the full source trace per fix.")


def closing_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(
        slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.7), align=PP_ALIGN.CENTER
    )
    line(
        frame,
        "v1.3.5 is ready to be tagged and published",
        34,
        WHITE,
        bold=True,
        first=True,
        after=14,
    )
    line(
        frame,
        "The fleet's onboarding, dispatch, and attestation paths are now proven at their original failure points, not just patched around.",
        20,
        PALE,
    )
    notes(
        slide,
        "Close with the practical consequence: these were mysteries at session start, now closed with regression tests.",
    )


for build in (
    title_slide,
    onboarding_slide,
    dispatch_slide,
    evidence_slide,
    boundary_slide,
    closing_slide,
):
    build()

prs.save(OUT)
print(OUT)
