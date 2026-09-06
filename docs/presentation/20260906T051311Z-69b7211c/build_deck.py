#!/usr/bin/env python3
"""Build the six-slide MAC v1.4.0 release deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
COMMIT = "69b7211c"
CAPTURED = "2026-09-06T05:13:11Z"
OUT = HERE / f"mac-capabilities-{COMMIT}.pptx"

INK = RGBColor(0x1B, 0x25, 0x30)
SLATE = RGBColor(0x2F, 0x3A, 0x45)
GREY = RGBColor(0x5A, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x98, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x6F, 0x9E)
AMBER = RGBColor(0xB5, 0x65, 0x1D)
SAND = RGBColor(0xF2, 0xC4, 0x8A)
PALE = RGBColor(0xDC, 0xE5, 0xEC)
FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def background(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.paragraphs[0].alignment = align
    return frame


def line(frame, text, size, color, *, bold=False, first=False, after=8):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(after)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def title_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.8))
    line(frame, "MAC", 54, WHITE, bold=True, first=True, after=8)
    line(frame, "v1.4.0", 32, PALE, after=16)
    line(
        frame,
        "The fleet's chat gateway is Hermes again: OpenClaw's onboarding was hardened, then its unfixable"
        " root cause forced a cutover.",
        18,
        MUTED,
    )
    footer = text_box(slide, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.8))
    line(footer, f"commit {COMMIT}   ·   captured {CAPTURED}", 16, SAND, bold=True, first=True)
    notes(
        slide,
        "This deck is pinned to the release commit. AUDIT.md traces each visible claim to the source tree or generated reference.",
    )


def hardening_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(
        frame,
        "OpenClaw was given an honest try first",
        33,
        INK,
        bold=True,
        first=True,
    )
    line(frame, "Three real reliability bugs, fixed before the root cause was found.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.05), Inches(11.3), Inches(4.65))
    for index, (head, detail) in enumerate(
        [
            (
                "Cron schedule collision",
                "Two hourly jobs sharing an identical schedule raced to open the same SQLite state file every hour.",
            ),
            (
                "Host-side flock mutex",
                "Every sandboxed OpenClaw CLI invocation now serializes through a single host-owned lock file.",
            ),
            (
                "Message-body encoding",
                "Delivery silently dropped the message body via an ignored --presentation field; fixed to encode through --message.",
            ),
        ]
    ):
        line(
            body, head, 21, BLUE if index % 2 == 0 else AMBER, bold=True, first=index == 0, after=2
        )
        line(body, detail, 17, SLATE, after=12)
    notes(slide, "Each point is a shipped fix, traced in AUDIT.md. None of them fixed the underlying problem.")


def root_cause_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(slide, Inches(0.9), Inches(0.6), Inches(11.5), Inches(1.3))
    line(frame, "The root cause was outside this repository", 33, WHITE, bold=True, first=True)
    body = text_box(slide, Inches(0.95), Inches(2.0), Inches(11.3), Inches(4.5))
    line(
        body,
        "OpenClaw's sandboxed state database sat on Docker Desktop's overlayfs.",
        22,
        PALE,
        bold=True,
        first=True,
        after=14,
    )
    line(
        body,
        "POSIX advisory locking is broken enough there that a fresh, empty SQLite WAL database hung"
        " indefinitely under a trivial write load — proven directly, not inferred.",
        19,
        MUTED,
        after=20,
    )
    line(
        body,
        "A high-severity bug report was filed upstream with the reproduction.",
        20,
        SAND,
        bold=True,
    )
    notes(
        slide,
        "This is the pivot slide: three prior fixes closed real bugs but couldn't touch a filesystem-level"
        " locking failure. See AUDIT.md for the upstream issue link.",
    )


def cutover_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(frame, "The fleet's chat gateway is Hermes again", 33, INK, bold=True, first=True)
    line(frame, "Not vendored, not pip-installed — upstream's own distribution model.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.05), Inches(11.3), Inches(4.65))
    for index, (head, detail) in enumerate(
        [
            (
                "Shell-installed, not vendored",
                "Upstream's official installer, run per node. The 444k-line vendored snapshot removed last month stays removed.",
            ),
            (
                "State migrated, not lost",
                "hermes claw migrate pulled each node's accumulated OpenClaw identity, memory, and skills across during cutover.",
            ),
            (
                "Repo-owned lifecycle, not manual",
                "deploy/hermes/install-hermes-gateway.sh (prepare/verify/finalize/withdraw) is now a first-class fleet-deploy option.",
            ),
        ]
    ):
        line(
            body, head, 21, BLUE if index % 2 == 0 else AMBER, bold=True, first=index == 0, after=2
        )
        line(body, detail, 17, SLATE, after=12)
    notes(slide, "Each point closes a gap between the manual cutover and durable, reproducible automation.")


def caught_live_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    frame = text_box(slide, Inches(0.85), Inches(0.45), Inches(11.6), Inches(1.35))
    line(frame, "A regression was caught live, not in CI", 33, INK, bold=True, first=True)
    line(frame, "Applying the new automation idempotently broke a working node.", 16, GREY)
    body = text_box(slide, Inches(0.95), Inches(2.1), Inches(11.3), Inches(4.5))
    line(
        body,
        "hermes config set model <value> silently replaces the entire nested model configuration"
        " object, discarding a working base_url and api_key.",
        22,
        AMBER,
        bold=True,
        first=True,
        after=16,
    )
    line(
        body,
        "Re-running the installer against an already-configured node broke model routing outright."
        " Caught by a live verification call before it reached a second node.",
        19,
        SLATE,
        after=16,
    )
    line(
        body,
        "Fixed with dotted-path keys (model.default / model.provider / model.base_url), which only"
        " touch the named sub-key.",
        20,
        BLUE,
        bold=True,
    )
    notes(slide, "This is the kind of bug a test suite doesn't catch: it only reproduces against real, prior state.")


def closing_slide() -> None:
    slide = prs.slides.add_slide(BLANK)
    background(slide, SLATE)
    frame = text_box(
        slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.7), align=PP_ALIGN.CENTER
    )
    line(
        frame,
        "v1.4.0 is ready to be tagged and published",
        34,
        WHITE,
        bold=True,
        first=True,
        after=14,
    )
    line(
        frame,
        "All three fleet nodes are live on Hermes, verified healthy, and durably managed —"
        " not a workaround, a real architectural correction.",
        20,
        PALE,
    )
    notes(
        slide,
        "Close with the practical consequence: chat-gateway reliability was a live production incident,"
        " now closed with root cause, fix, and durable automation.",
    )


for build in (
    title_slide,
    hardening_slide,
    root_cause_slide,
    cutover_slide,
    caught_live_slide,
    closing_slide,
):
    build()

prs.save(OUT)
print(OUT)
