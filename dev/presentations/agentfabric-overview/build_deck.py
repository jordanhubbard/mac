#!/usr/bin/env python3
"""Build the AgentFabric overview deck with python-pptx.

Deterministic, no LLM calls, no plugin runtime. The primitive shape vocabulary
(`shape`, `text`, `pill`, `chip`, `line`, `arrow`, `footer`, `title`, `wash`,
`labeled_card`) is inherited from the authoring package this directory was
seeded from; every diagram here is drawn with native PowerPoint shapes so the
mechanism stays readable, editable, and diffable after a Google Slides import.

`deck-specification.md` is the narrative authority for this file and
`source-notes.md` is its factual ledger. Slide order follows the specification:
marketing first, then reused technology, then actors and roles.

Usage:
    python3 build_deck.py
Environment overrides:
    AGENTFABRIC_DECK_SOURCE, AGENTFABRIC_REPO, AGENTFABRIC_DECK_OUTPUT
"""

from __future__ import annotations

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
SOURCE = Path(os.environ.get("AGENTFABRIC_DECK_SOURCE", str(HERE)))
REPO = Path(os.environ.get("AGENTFABRIC_REPO", str(SOURCE.parents[2])))


def _obj_dir() -> Path:
    return Path(os.environ.get("OBJ_DIR") or (REPO / "_build"))


OUT = Path(
    os.environ.get("AGENTFABRIC_DECK_OUTPUT")
    or (REPO / "docs" / "presentations" / "agentfabric-overview" / "agentfabric-overview.pptx")
)

W, H, TOTAL = 1280, 720, 20
PX = 9525  # EMU per px at 96dpi.
BRAND = "AGENTFABRIC"

C = {
    "ink": "#101317",
    "panel": "#23282F",
    "steel": "#65707C",
    "fog": "#EEF1F3",
    "white": "#FFFFFF",
    "orange": "#FF6B35",
    "orange2": "#FF9B66",
    "blue": "#72B7D6",
    "green": "#76B900",  # NVIDIA green: reserved for reused NVIDIA technology.
    "green2": "#7BC6A4",
    "red": "#F47C7C",
    "line": "#D8DEE3",
    "muted": "#AAB3BC",
    "code": "#1A1F26",
    "codeline": "#333B45",
}
FONT = "Helvetica Neue"
MONO = "Courier New"

_GEOMETRY = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
}
_RADIUS = {"rounded-xl": 0.14, "rounded-lg": 0.09, "rounded-full": 0.5}
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _px(v: float) -> Emu:
    return Emu(round(v * PX))


def _set_alpha(fill, opacity: float) -> None:
    solid_fill = fill.fore_color._xFill
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = solid_fill.makeelement(qn("a:alpha"), {"val": str(round(opacity * 100000))})
    srgb.append(alpha)


def shape(slide, geometry, x, y, w, h, fill, **opts):
    sp = slide.shapes.add_shape(_GEOMETRY[geometry], _px(x), _px(y), _px(w), _px(h))
    sp.shadow.inherit = False
    if fill == "none":
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        opacity = opts.get("opacity")
        if opacity is not None:
            _set_alpha(sp.fill, opacity)
    stroke = opts.get("stroke", "none")
    if stroke in (None, "none"):
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(stroke)
        sp.line.width = Pt(opts.get("strokeWidth") or 1)
    radius_key = opts.get("radius")
    if radius_key and geometry == "roundRect":
        sp.adjustments[0] = _RADIUS.get(radius_key, 0.05)
    return sp


def text(slide, value, x, y, w, h, size=28, color=C["ink"], bold=False, **opts):
    box = slide.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.text = value
    align = _ALIGN.get(opts.get("align"))
    font_name = opts.get("font", FONT)
    italic = opts.get("italic", False)
    for para in tf.paragraphs:
        if align is not None:
            para.alignment = align
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = _rgb(color)
    if opts.get("fit"):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return box


def mono(slide, value, x, y, w, h, size=12, color=C["fog"], bold=False):
    return text(slide, value, x, y, w, h, size, color, bold, font=MONO)


def line(slide, x, y, w, h=4, color=C["orange"]):
    return shape(slide, "rect", x, y, w, h, color)


def dot(slide, x, y, d, color, stroke="none"):
    return shape(
        slide,
        "ellipse",
        x,
        y,
        d,
        d,
        color,
        stroke=stroke,
        strokeWidth=0 if stroke == "none" else 2,
    )


def pill(slide, label, x, y, w, fill=C["orange"], color=C["white"], size=13):
    shape(slide, "roundRect", x, y, w, 32, fill, radius="rounded-full")
    text(slide, label, x + 10, y + 6, w - 20, 20, size, color, True, align="center", fit=True)


def chip(slide, label, x, y, w, fill=C["panel"], color=C["fog"], size=13):
    shape(slide, "roundRect", x, y, w, 34, fill, radius="rounded-lg")
    text(slide, label, x + 12, y + 7, w - 24, 20, size, color, True, fit=True)


def labeled_card(
    slide,
    x,
    y,
    w,
    h,
    fill,
    heading,
    body,
    heading_color,
    body_color,
    heading_size=20,
    body_size=15,
    **opts,
):
    """A rounded cell whose copy lives in the shape text frame, so Google cannot overflow it."""
    sp = shape(slide, "roundRect", x, y, w, h, fill, radius="rounded-xl", **opts)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = _px(18)
    tf.margin_top = tf.margin_bottom = _px(18)
    lines = [heading, *body.split("\n")]
    for i, entry in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.text = entry
        if i == 0:
            para.space_after = Pt(10)
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(heading_size if i == 0 else body_size)
            run.font.bold = True
            run.font.color.rgb = _rgb(heading_color if i == 0 else body_color)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return sp


def caption_bar(slide, value, y=648, fill="#FFF0E8", color=C["ink"]):
    shape(slide, "roundRect", 70, y, 1140, 36, fill, radius="rounded-full")
    text(slide, value, 90, y + 8, 1100, 22, 13, color, True, align="center", fit=True)


def wash(slide, x=0, y=0, w=W, h=H, opacity=0.35):
    return shape(slide, "rect", x, y, w, h, C["ink"], opacity=opacity)


def title(slide, value, sub, dark, n, tag=None):
    # Reserved header band (y=28-252). Title boxes are taller than the type they
    # hold so Google's Arial substitution cannot paint into the subtitle.
    ink = C["white"] if dark else C["ink"]
    muted = C["muted"] if dark else C["steel"]
    pill(
        slide,
        BRAND,
        70,
        28,
        148,
        C["white"] if dark else C["ink"],
        C["ink"] if dark else C["white"],
    )
    if tag:
        tag_w = max(188, 32 + len(tag) * 9)
        pill(slide, tag, 230, 28, tag_w, C["orange"], C["ink"], 12)
    text(slide, str(n).zfill(2), 1170, 32, 40, 22, 13, ink, True, align="right")
    long = len(value) > 46
    t_size = 32 if long else 38
    text(slide, value, 70, 78, 1140, 84, t_size, ink, True, fit=True)
    if sub:
        text(slide, sub, 72, 168, 1100, 68, 17, muted, False, fit=True)
    line(slide, 70, 248, 88, 6, C["orange"])


def footer(slide, n, dark=False):
    # Footer lives below y=696 so bottom captions ending at y=684 cannot collide.
    text(slide, BRAND, 70, 698, 170, 16, 11, C["muted"] if dark else C["steel"], True)
    line(slide, 1080, 704, 120, 3, C["panel"] if dark else C["line"])
    line(slide, 1080, 704, max(8, 120 * n / TOTAL), 3, C["orange"])


def arrow(slide, x1, y, x2, color=C["orange"]):
    line(slide, x1, y, max(4, x2 - x1 - 13), 4, color)
    shape(slide, "chevron", x2 - 19, y - 8, 19, 20, color)


def down_arrow(slide, x, y1, y2, color=C["orange"]):
    shape(slide, "rect", x, y1, 4, max(4, y2 - y1 - 12), color)
    sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, _px(x - 8), _px(y2 - 14), _px(20), _px(18))
    sp.shadow.inherit = False
    sp.rotation = 90
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(color)
    sp.line.fill.background()
    return sp


def notes(slide, lines):
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def section(slide, kicker, headline, body, n):
    """Full-bleed divider: one kicker, one headline, one sentence of scope."""
    pill(slide, BRAND, 70, 48, 148, C["white"], C["ink"])
    text(slide, kicker, 70, 168, 700, 30, 15, C["orange2"], True)
    text(slide, headline, 70, 208, 900, 150, 46, C["white"], True, fit=True)
    line(slide, 74, 384, 120, 8, C["orange"])
    text(slide, body, 74, 414, 760, 110, 20, C["fog"], False, fit=True)
    footer(slide, n, True)


def stage_flow(slide, labels, y, x0=70, width=1140, height=44, dark=False, accent=None):
    """A left-to-right chevron flow. Diagram, not prose: labels are one or two words."""
    accent = accent or C["orange"]
    n = len(labels)
    step = width / n
    body_w = step - 12
    for i, label in enumerate(labels):
        x = x0 + i * step
        fill = accent if i == 0 else (C["panel"] if not dark else C["white"])
        color = C["ink"] if i == 0 or dark else C["white"]
        shape(slide, "chevron", x, y, body_w, height, fill)
        text(
            slide,
            label,
            x + 8,
            y + (height - 22) / 2,
            body_w - 30,
            22,
            13,
            color,
            True,
            align="center",
            fit=True,
        )


def build() -> Path:
    p = Presentation()
    p.slide_width = _px(W)
    p.slide_height = _px(H)
    blank = p.slide_layouts[6]

    def add_slide(bg=C["white"]):
        s = p.slides.add_slide(blank)
        shape(s, "rect", 0, 0, W, H, bg)
        return s

    # ------------------------------------------------------------------ part 0
    # 01 - cover
    s = add_slide(C["ink"])
    # Hero: a hub with a fleet around it, drawn natively. No raster, no text in art.
    shape(s, "ellipse", 900, 250, 220, 220, C["panel"], stroke=C["orange"], strokeWidth=2)
    text(s, "HUB", 940, 340, 140, 30, 18, C["orange"], True, align="center")
    fleet = [(760, 170), (1140, 170), (700, 400), (1180, 400), (760, 560), (1140, 560)]
    for fx, fy in fleet:
        dot(s, fx, fy, 56, C["panel"], C["blue"])
        arrow(s, min(fx + 56, 1010), fy + 26, max(fx + 70, 902), C["steel"])
    pill(s, BRAND, 70, 55, 148, C["white"], C["ink"])
    text(
        s,
        "One control plane\nfor a fleet of\nAI agents.",
        70,
        140,
        560,
        220,
        44,
        C["white"],
        True,
        fit=True,
    )
    line(s, 74, 402, 112, 7, C["orange"])
    text(
        s,
        "AgentFabric turns a conversation into durable work, dispatches it across a heterogeneous"
        " fleet of agents and models, and keeps the receipts.",
        74,
        428,
        520,
        130,
        19,
        C["fog"],
        False,
        fit=True,
    )
    text(
        s,
        "MULTI-AGENT, MULTI-MODEL ORCHESTRATION AT SCALE",
        74,
        580,
        560,
        28,
        14,
        C["orange2"],
        True,
        fit=True,
    )
    notes(
        s,
        [
            "Audience assumption: highly technical, but new both to AgentFabric and to fleet-scale agent orchestration. Nothing in this deck assumes prior product knowledge.",
            "Order of the deck: what the system does (1-5), what it reuses and why (6-10), then who the actors are and how the machine actually runs (11-19).",
            "Every quantified figure traces to source-notes.md, which cites the control-plane tree it was measured from. No modeled ROI or productivity numbers appear anywhere in this deck.",
        ],
    )

    # 02 - what it does
    s = add_slide(C["white"])
    title(
        s,
        "Ask once. The fabric carries it to production.",
        "A request becomes a durable task, the task is leased to an agent that can actually do it, the"
        " result is reviewed by a different agent, and every step leaves evidence.",
        False,
        2,
        tag="WHAT IT DOES",
    )
    stage_flow(
        s,
        ["ASK", "DURABLE TASK", "LEASE + DISPATCH", "EXECUTE", "INDEPENDENT REVIEW", "PUBLISH"],
        300,
    )
    cards2 = [
        ("NOTHING IS LOST", "A task outlives the\nchat, the agent, and\nthe machine.", C["orange"]),
        ("NOTHING IS BLIND", "Cost, timing, and\nprovenance are recorded\nper step.", C["blue"]),
        (
            "NOTHING SELF-APPROVES",
            "The agent that did the\nwork does not accept\nthe work.",
            C["green2"],
        ),
    ]
    for i, (head, body, color) in enumerate(cards2):
        x = 70 + i * 386
        labeled_card(
            s,
            x,
            396,
            368,
            200,
            C["fog"],
            head,
            body,
            C["ink"],
            C["steel"],
            19,
            15,
            stroke=color,
            strokeWidth=2,
        )
    caption_bar(s, "The unit of truth is a task in a ledger, not a message in a transcript.")
    footer(s, 2)
    notes(
        s,
        [
            "This is the marketing frame, and it is also literally the pipeline: ask, durable task, lease, execute, independent review, publish.",
            "'Independent review' is a separate agent invocation under control-plane-owned gates, not the executing agent grading itself.",
            "Mechanism for each of these three claims appears later: the ledger and its state machine on slide 13, measurement on slide 17, and the review gates on slides 12-13.",
        ],
    )

    # 03 - the problem
    s = add_slide(C["fog"])
    title(
        s,
        "A chat window is not an operating model.",
        "One human supervising one agent works. The same pattern at fleet scale loses work, repeats"
        " work, and cannot answer what any of it cost.",
        False,
        3,
        tag="THE PROBLEM",
    )
    shape(
        s,
        "roundRect",
        70,
        300,
        520,
        296,
        C["white"],
        radius="rounded-xl",
        stroke=C["line"],
        strokeWidth=1,
    )
    text(s, "ONE AGENT", 100, 326, 300, 26, 15, C["steel"], True)
    dot(s, 100, 380, 64, C["orange"])
    text(s, "human", 96, 456, 72, 22, 13, C["steel"], True, align="center")
    arrow(s, 168, 410, 300)
    dot(s, 300, 380, 64, C["panel"])
    text(s, "agent", 296, 456, 72, 22, 13, C["steel"], True, align="center")
    arrow(s, 368, 410, 500)
    dot(s, 470, 380, 64, C["blue"])
    text(s, "result", 466, 456, 72, 22, 13, C["steel"], True, align="center")
    text(
        s,
        "Supervision is the human. It works, and it does not scale.",
        100,
        516,
        460,
        56,
        15,
        C["ink"],
        True,
        fit=True,
    )

    shape(s, "roundRect", 620, 300, 590, 296, C["ink"], radius="rounded-xl")
    text(s, "FORTY AGENTS", 650, 326, 340, 26, 15, C["orange2"], True)
    for i in range(8):
        for j in range(4):
            dot(s, 654 + i * 62, 366 + j * 44, 26, C["panel"] if (i + j) % 3 else C["red"])
    text(
        s,
        "Same pattern, no system of record: duplicated effort, stranded work,\n"
        'unattributed spend, and no answer to "who approved this?"',
        650,
        526,
        540,
        56,
        15,
        C["fog"],
        True,
        fit=True,
    )
    footer(s, 3)
    notes(
        s,
        [
            "The failure is structural, not a tooling gap: conversation is ephemeral state, and a fleet needs durable state with ownership, ordering, and receipts.",
            "The red nodes on the right are illustrative of failure classes (stranded, duplicated, unattributed), not a measured failure rate. Do not attach a percentage to this slide.",
            "Real ledger census figures, when quoted, come from source-notes.md and carry the date they were measured.",
        ],
    )

    # 04 - what you get
    s = add_slide(C["ink"])
    title(
        s,
        "Four properties the fabric guarantees.",
        "These are the reasons to put a control plane underneath agents at all, rather than scaling"
        " prompts and hoping.",
        True,
        4,
        tag="WHY IT MATTERS",
    )
    props = [
        (
            "DURABLE TRUTH",
            "Tasks, leases, reviews,\nsecrets, and audit trails\nsurvive every restart.",
            C["orange"],
        ),
        (
            "HONEST HETEROGENEITY",
            "Different machines,\nOSes, agents, and models\nare modelled, not hidden.",
            C["blue"],
        ),
        (
            "EVIDENCE BY DEFAULT",
            "Build, test, review, and\npublication decisions are\nnamed and stored.",
            C["green2"],
        ),
        (
            "SPEND VISIBILITY",
            "Model calls are recorded\nas route events and\npriced at read time.",
            C["orange2"],
        ),
    ]
    for i, (head, body, color) in enumerate(props):
        x = 70 + i * 290
        labeled_card(
            s,
            x,
            300,
            272,
            240,
            C["panel"],
            head,
            body,
            color,
            C["fog"],
            18,
            14,
            stroke=color,
            strokeWidth=2,
        )
    text(
        s,
        "A fleet you cannot audit is a fleet you cannot trust with a repository.",
        70,
        576,
        1140,
        40,
        22,
        C["white"],
        True,
        align="center",
        fit=True,
    )
    footer(s, 4, True)
    notes(
        s,
        [
            "Durable truth: the control plane owns tasks, leases, routing, reviews, evidence, secret handles, runtime manifests, rollout state, and audit trails.",
            "Honest heterogeneity: node classes differ by design (host installs versus containerized execution), and the dispatcher matches work to capability rather than pretending the fleet is uniform.",
            "Spend visibility: usage is recorded where routing happens, one event per model call, and cost is priced at read time - so re-pricing history does not require re-running work. Enforcing metering at the router (rather than trusting the caller to report usage) is an accepted-in-principle but still proposed decision; today a streamed route that did not request usage is recorded as a coverage gap.",
        ],
    )

    # 05 - end to end
    s = add_slide(C["white"])
    title(
        s,
        "One request, end to end.",
        "The same path every piece of work takes. Read it left to right; the annotations under each"
        " hand-off are what the control plane writes down.",
        False,
        5,
        tag="THE LOOP",
    )
    lanes = [
        ("HUMAN", "asks in chat", C["orange"]),
        ("GATEWAY", "conversation, memory", C["blue"]),
        ("HUB", "task + lease", C["ink"]),
        ("WORKER", "sandboxed execution", C["panel"]),
        ("REVIEWER", "independent decision", C["green2"]),
        ("PUBLISH", "merge queue, receipts", C["orange2"]),
    ]
    for i, (head, sub, color) in enumerate(lanes):
        x = 70 + i * 192
        shape(
            s,
            "roundRect",
            x,
            300,
            172,
            96,
            C["fog"],
            radius="rounded-xl",
            stroke=color,
            strokeWidth=2,
        )
        text(s, head, x + 14, 322, 144, 24, 15, C["ink"], True, align="center", fit=True)
        text(s, sub, x + 14, 350, 144, 34, 12, C["steel"], False, align="center", fit=True)
        if i < len(lanes) - 1:
            arrow(s, x + 172, 346, x + 192)
        down_arrow(s, x + 84, 400, 452, C["line"])
    records = [
        "intent",
        "identity",
        "ledger row,\nstate, owner",
        "action events,\ntoken spend",
        "named gate\ndecision",
        "provenance,\nevidence pointer",
    ]
    for i, record in enumerate(records):
        x = 70 + i * 192
        shape(s, "roundRect", x, 452, 172, 96, C["ink"], radius="rounded-lg")
        text(s, record, x + 14, 474, 144, 60, 13, C["fog"], True, align="center", fit=True)
    text(s, "WHAT GETS WRITTEN DOWN", 70, 566, 400, 26, 14, C["steel"], True)
    caption_bar(
        s, "Every arrow crosses a boundary the control plane owns; nothing advances on trust alone."
    )
    footer(s, 5)
    notes(
        s,
        [
            "This slide is the whole system in one picture; slides 12 through 17 expand each lane.",
            "The gateway owns conversation, personality, and memory. It deliberately does not own operational truth - that separation is the core architectural decision.",
            "'Named gate decision' is literal: gates return a named decision rather than a boolean, so a refusal is explainable after the fact.",
        ],
    )

    # ------------------------------------------------------------------ part 1
    # 06 - divider: reuse
    s = add_slide(C["ink"])
    section(
        s,
        "PART ONE",
        "Built on the stack, not instead of it.",
        "Isolation, observability, inference, containers, scheduling, and coding agents already exist"
        " and are maintained by people whose full-time job they are. AgentFabric integrates them and"
        " spends its own code on the one thing none of them provide: durable operational truth.",
        6,
    )
    notes(
        s,
        [
            "Framing for the engineering audience: the interesting question is not what was written, it is what was deliberately not written.",
            "Each reuse claim on the next three slides names the project and the boundary. Lineage (learned from) is stated separately from integration (depends on) on purpose.",
        ],
    )

    # 07 - NVIDIA technology
    s = add_slide(C["white"])
    title(
        s,
        "Where AgentFabric leverages NVIDIA technology.",
        "Four NVIDIA projects carry four hard problems: confinement, observability, elastic GPU"
        " capacity, and a reference conversational runtime.",
        False,
        7,
        tag="NVIDIA STACK",
    )
    shape(s, "roundRect", 70, 296, 1140, 74, C["ink"], radius="rounded-xl")
    text(s, "AGENTFABRIC CONTROL PLANE", 90, 314, 500, 30, 18, C["white"], True)
    text(
        s,
        "tasks · leases · routing · reviews · evidence · audit",
        700,
        320,
        490,
        26,
        14,
        C["muted"],
        False,
        align="right",
        fit=True,
    )
    nv = [
        (
            "NVIDIA OpenShell",
            "EXECUTION SECURITY",
            "Landlock filesystem policy,\nseccomp syscall filter,\ndeny-by-default L7 egress.\nOne guardrail authority.",
        ),
        (
            "NVIDIA NeMo Relay",
            "OBSERVABILITY",
            "Request, task, tool, and\nmodel activity mapped into\nRelay scopes when the\n`relay` extra is enabled.",
        ),
        (
            "NVIDIA HGX",
            "ELASTIC CAPACITY",
            "Bounded, receipt-bearing\nautoscaling of provider\nsessions from durable\nprovisioning demand.",
        ),
        (
            "NVIDIA NemoClaw",
            "REFERENCE INTEGRATION",
            "Compatibility and design\nreference for the\nconversational agent\nboundary.",
        ),
    ]
    for i, (name, role, body) in enumerate(nv):
        x = 70 + i * 290
        shape(
            s,
            "roundRect",
            x,
            404,
            272,
            214,
            C["fog"],
            radius="rounded-xl",
            stroke=C["green"],
            strokeWidth=2,
        )
        down_arrow(s, x + 134, 370, 404, C["green"])
        text(s, name, x + 16, 424, 240, 30, 16, C["ink"], True, align="center", fit=True)
        text(s, role, x + 16, 456, 240, 22, 11, C["green"], True, align="center", fit=True)
        line(s, x + 106, 484, 60, 3, C["green"])
        text(s, body, x + 16, 500, 240, 104, 12, C["steel"], False, align="center", fit=True)
    footer(s, 7)
    notes(
        s,
        [
            "OpenShell is the security boundary: process trees, filesystem and network policy, sandbox lifecycle, and normalized action-event collection integrate with it rather than reimplementing isolation. The design goal is exactly one guardrail authority, the policy file, instead of two competing ones.",
            "NeMo Relay is optional and enabled through a packaging extra; observability is mapped into it rather than invented.",
            "HGX capacity is operator-side and bounded: read-only status and plan commands are separate from an explicit execute, and provider work never runs on a dispatcher or HTTP thread.",
            "NemoClaw is reference material for the conversational boundary, not the deployed implementation of it. That distinction is deliberate and is recorded in source-notes.md.",
        ],
    )

    # 08 - OSS technology
    s = add_slide(C["fog"])
    title(
        s,
        "Where AgentFabric leverages open source.",
        "Bullets, here, on purpose: this is the list of projects doing load-bearing work, grouped by"
        " the problem each one already solved.",
        False,
        8,
        tag="OSS STACK",
    )
    groups = [
        (
            "STATE",
            C["orange"],
            ["PostgreSQL", "SQLite (local)", "Alembic-style\nversioned migrations"],
        ),
        ("SERVICE", C["blue"], ["FastAPI", "Uvicorn", "Pydantic", "httpx"]),
        (
            "EXECUTION",
            C["green2"],
            ["Docker Engine / Moby", "Kubernetes", "OpenClaw", "Codex CLI · OpenCode"],
        ),
        ("PROTOCOL", C["orange2"], ["ACP", "A2A agent cards", "MCP", "OCSF event streams"]),
    ]
    for i, (head, color, items) in enumerate(groups):
        x = 70 + i * 290
        shape(
            s,
            "roundRect",
            x,
            300,
            272,
            320,
            C["white"],
            radius="rounded-xl",
            stroke=color,
            strokeWidth=2,
        )
        text(s, head, x + 20, 322, 232, 28, 16, C["ink"], True, fit=True)
        line(s, x + 20, 356, 46, 4, color)
        for j, item in enumerate(items):
            dot(s, x + 22, 384 + j * 58, 10, color)
            text(s, item, x + 44, 376 + j * 58, 208, 48, 14, C["steel"], True, fit=True)
    caption_bar(
        s,
        "Every box above is a dependency, not a fork. The control plane owns none of these problems.",
    )
    footer(s, 8)
    notes(
        s,
        [
            "PostgreSQL is the authority for fleet deployments; the local development path is SQLite, and schema changes require versioned migrations with fail-closed authority.",
            "Coding executors are plural and pluggable - the route ladder on slide 16 orders them - because no single coding agent is uniformly best or uniformly available.",
            "Protocols are implemented specifications, not aspirations: ACP and A2A endpoints and agent cards are served by the hub, MCP is a client of the same API rather than a second implementation, and OCSF is the event vocabulary produced under sandboxed execution.",
            "Hermes Agent is lineage rather than dependency: the fleet learned from its loop, gateways, and skills, and the in-tree snapshot was removed. Say 'learned from', not 'built on', for that one.",
        ],
    )

    # 09 - what AgentFabric adds
    s = add_slide(C["ink"])
    title(
        s,
        "What AgentFabric adds that nothing above provides.",
        "Reuse is the default. These five mechanisms are the parts that had to be built, because"
        " they are the parts that make agent output trustworthy.",
        True,
        9,
        tag="UNIQUE POWERS",
    )
    powers = [
        (
            "DURABLE LEDGER",
            "A task is a state machine\nwith an owner, a lease,\ndependencies, and history.",
        ),
        (
            "NAMED GATES",
            "A gate returns a named\ndecision, not a boolean.\nRefusals are explainable.",
        ),
        (
            "ROUTE LADDER",
            "Ordered, capability-matched\nselection across coding\nagents and models.",
        ),
        (
            "EVIDENCE CLOSURE",
            "Build, test, review, and\npublication artifacts are\npointers, kept with the task.",
        ),
        (
            "BREAK-GLASS",
            "Recovery is a granted,\nlistable, revocable\nauthorization with a reason.",
        ),
    ]
    for i, (head, body) in enumerate(powers):
        x = 70 + i * 231
        shape(
            s,
            "roundRect",
            x,
            300,
            213,
            250,
            C["panel"],
            radius="rounded-xl",
            stroke=C["orange"],
            strokeWidth=2,
        )
        text(s, str(i + 1).zfill(2), x + 20, 320, 60, 26, 14, C["orange"], True)
        text(s, head, x + 20, 352, 176, 56, 16, C["white"], True, fit=True)
        line(s, x + 20, 414, 44, 3, C["orange"])
        text(s, body, x + 20, 432, 176, 100, 13, C["fog"], False, fit=True)
    text(
        s,
        "Agent output is a candidate. These five mechanisms are how it earns the right to land.",
        70,
        580,
        1140,
        40,
        21,
        C["orange2"],
        True,
        align="center",
        fit=True,
    )
    footer(s, 9, True)
    notes(
        s,
        [
            "This is the anti-reinvention slide's counterweight: the deck is explicit that reuse is preferred, so it must be equally explicit about the parts that are genuinely ours.",
            "Leases are what make concurrency safe: a claim is time-bounded and recoverable, so a dead worker does not strand work forever - there are explicit recovery verbs for stranded and stalled work.",
            "Break-glass exists because recovery paths that are undocumented become undocumented root access. Granting, listing, and revoking are all first-class and reason-bearing.",
            "Precision on named gates: the review gate returns a named decision today; making that the contract for every gate is an accepted-in-principle proposal rather than a finished refactor. Say it that way if asked.",
        ],
    )

    # 10 - trust boundary
    s = add_slide(C["white"])
    title(
        s,
        "The trust boundary, drawn explicitly.",
        "Agents run with wide latitude inside a narrow box. The box - not the agent's good"
        " judgement - is what makes autonomous execution acceptable.",
        False,
        10,
        tag="CONTAINMENT",
    )
    shape(
        s,
        "roundRect",
        70,
        296,
        700,
        324,
        C["fog"],
        radius="rounded-xl",
        stroke=C["orange"],
        strokeWidth=2,
    )
    text(s, "INSIDE THE SANDBOX", 96, 316, 400, 26, 15, C["orange"], True)
    inside = [
        "Agent process tree, owned and reaped",
        "Read-only and read-write paths, allow-listed",
        "Syscall filter; never runs as root",
        "Declared egress hosts, per project and per task",
        "Secret handles resolved at use, never printed",
    ]
    for i, item in enumerate(inside):
        dot(s, 96, 360 + i * 48, 12, C["orange"])
        text(s, item, 122, 352 + i * 48, 620, 36, 15, C["ink"], True, fit=True)
    shape(s, "roundRect", 800, 296, 410, 324, C["ink"], radius="rounded-xl")
    text(s, "NEVER CROSSES", 826, 316, 340, 26, 15, C["red"], True)
    outside = [
        "Undeclared network destinations",
        "Raw credential values",
        "Another project's repository",
        "Self-approval of its own work",
    ]
    for i, item in enumerate(outside):
        shape(s, "rect", 826, 366 + i * 58, 26, 4, C["red"])
        text(s, item, 866, 352 + i * 58, 320, 44, 15, C["fog"], True, fit=True)
    caption_bar(
        s,
        "Deny-by-default egress and secrets-as-handles are policy, enforced by the sandbox, not conventions.",
    )
    footer(s, 10)
    notes(
        s,
        [
            "The containment model is OpenShell's, configured by a control-plane-authored policy: filesystem via Landlock, syscalls via seccomp, network via an L7 proxy that is deny-by-default.",
            "Secrets are handles: work references a secret by name and the value is resolved at use, so a transcript, a log, or an evidence bundle never carries the value.",
            "'Never self-approves' is an authority boundary rather than a sandbox control, and it is the one item on the right that the control plane enforces itself.",
        ],
    )

    # ------------------------------------------------------------------ part 2
    # 11 - divider: actors
    s = add_slide(C["ink"])
    section(
        s,
        "PART TWO",
        "Actors, roles, and who is allowed to decide.",
        "The rest of the deck is the machine itself: the cast of actors, the life of one task, how"
        " agents coordinate without a switchboard, how the fleet is shaped, and how models and money"
        " are routed and measured.",
        11,
    )
    notes(
        s,
        [
            "Transition point. A decision audience can stop here; an engineering audience should stay for slides 12 through 17.",
            "The organizing idea for this half: every actor has exactly one kind of authority, and no actor holds two that would let it grade its own work.",
        ],
    )

    # 12 - the cast
    s = add_slide(C["white"])
    title(
        s,
        "The cast, and the one authority each actor holds.",
        "Roles are separated so that the ability to do work, the ability to accept work, and the"
        " ability to change the rules never live in the same place.",
        False,
        12,
        tag="ACTORS",
    )
    cast = [
        ("REQUESTER", "states intent", C["orange"]),
        ("GATEWAY AGENT", "conversation + memory", C["blue"]),
        ("HUB", "durable truth", C["ink"]),
        ("DISPATCHER", "matches work to capability", C["panel"]),
        ("WORKER", "runs one leased task", C["steel"]),
        ("CODING EXECUTOR", "edits the repository", C["orange2"]),
        ("REVIEWER AGENT", "returns a named decision", C["green2"]),
        ("OPERATOR", "grants, revokes, recovers", C["red"]),
    ]
    for i, (head, role, color) in enumerate(cast):
        col, row = i % 4, i // 4
        x = 70 + col * 290
        y = 300 + row * 168
        shape(
            s,
            "roundRect",
            x,
            y,
            272,
            148,
            C["fog"],
            radius="rounded-xl",
            stroke=color,
            strokeWidth=2,
        )
        shape(s, "rect", x + 20, y + 24, 40, 4, color)
        text(s, head, x + 20, y + 44, 232, 30, 17, C["ink"], True, fit=True)
        text(s, role, x + 20, y + 84, 232, 44, 14, C["steel"], True, fit=True)
    footer(s, 12)
    notes(
        s,
        [
            "The gateway agent is a stock conversational runtime under a control-plane-authored sandbox policy; it owns personality, memory, and channels, and holds no operational authority.",
            "Worker and coding executor are separate actors on purpose: the worker is the leased, supervised process; the executor is whichever coding agent the route ladder selected for it.",
            "The reviewer is a different invocation from the one that produced the work. The operator is the only actor that can grant recovery authority, and every grant is recorded with its reason.",
            "Missing from this list, deliberately: any actor that can both execute and accept. That combination does not exist in the model.",
        ],
    )

    # 13 - life of a task
    s = add_slide(C["fog"])
    title(
        s,
        "The life of one task.",
        "A task is a state machine, not a status string. Terminal states are reached by decision,"
        " and the path is stored with the task.",
        False,
        13,
        tag="LIFECYCLE",
    )
    stage_flow(
        s, ["OPEN", "WAITING", "CLAIMED", "RUNNING", "NEEDS REVIEW", "REVIEWING", "COMPLETED"], 300
    )
    gates = [
        ("DEPENDENCY", "unfinished dependency\nholds the task"),
        ("LEASE", "one owner, time-bounded,\nrecoverable"),
        ("EVIDENCE", "build and test output\nbefore review"),
        ("ACCEPTANCE", "an independent named\ndecision"),
    ]
    for i, (head, body) in enumerate(gates):
        x = 70 + i * 290
        shape(
            s,
            "roundRect",
            x,
            396,
            272,
            132,
            C["white"],
            radius="rounded-xl",
            stroke=C["orange"],
            strokeWidth=2,
        )
        text(s, f"GATE {i + 1}", x + 20, 414, 232, 22, 12, C["orange"], True)
        text(s, head, x + 20, 438, 232, 26, 16, C["ink"], True, fit=True)
        text(s, body, x + 20, 468, 232, 46, 13, C["steel"], False, fit=True)
        down_arrow(s, x + 134, 348, 396, C["line"])
    others = ["BLOCKED", "NEEDS INPUT", "STOPPED", "FAILED", "CANCELLED"]
    text(s, "ALSO TERMINAL OR HELD:", 70, 552, 290, 24, 13, C["steel"], True)
    for i, label in enumerate(others):
        chip(s, label, 380 + i * 166, 546, 158, C["ink"], C["fog"], 12)
    caption_bar(s, "Failure keeps its evidence. A failed task is a record, not a deletion.")
    footer(s, 13)
    notes(
        s,
        [
            "The state set is larger than the happy path shown here; the held and terminal states are listed along the bottom so the diagram stays readable while remaining complete about the vocabulary.",
            "Holds are not states: a staged task carries a hold flag and is released, which is why 'dispatchable' is a function of dependencies, capability, lease, and project pause rather than a single field.",
            "Recovery verbs exist for stranded and stalled work precisely because leases expire while a real process may or may not still be alive.",
        ],
    )

    # 14 - coordination
    s = add_slide(C["ink"])
    title(
        s,
        "Coordination is a town square, not a switchboard.",
        "Agents broadcast on a shared bus and act on what they hear. The hub does not sit in the"
        " middle of every conversation - it sits underneath the durable consequences.",
        True,
        14,
        tag="COORDINATION",
    )
    shape(s, "roundRect", 430, 396, 420, 76, C["orange"], radius="rounded-full")
    text(
        s, "AGENTBUS  ·  BROADCAST", 450, 418, 380, 32, 20, C["ink"], True, align="center", fit=True
    )
    participants = ["GATEWAY", "WORKER", "REVIEWER", "OPERATOR", "DISPATCHER"]
    for i, label in enumerate(participants):
        x = 70 + i * 231
        shape(
            s,
            "roundRect",
            x,
            300,
            213,
            62,
            C["panel"],
            radius="rounded-lg",
            stroke=C["blue"],
            strokeWidth=1,
        )
        text(s, label, x + 16, 318, 181, 26, 14, C["fog"], True, align="center", fit=True)
        down_arrow(s, x + 106, 366, 396, C["blue"])
    verbs = ["stand down", "abort", "pause", "resume", "status"]
    text(s, "LIFECYCLE VERBS ON THE BUS", 70, 512, 400, 26, 14, C["muted"], True)
    for i, label in enumerate(verbs):
        chip(s, label, 70 + i * 231, 548, 213, C["white"], C["ink"], 13)
    text(
        s,
        "Broadcast for intent and interruption. Ledger for consequence.",
        70,
        604,
        1140,
        34,
        18,
        C["orange2"],
        True,
        align="center",
        fit=True,
    )
    footer(s, 14, True)
    notes(
        s,
        [
            "The bus is a broadcast medium: a message is heard by the fleet rather than addressed through a central router, which is what makes an interruption like 'stand down' cheap and immediate.",
            "The durable consequence of anything heard on the bus still lands in the ledger. The bus is not the system of record and must not be described as one.",
            "Operational learning is recorded as secret-free memories, so repeated failures against the same credential pattern change future routing instead of being retried blindly.",
        ],
    )

    # 15 - fleet shape
    s = add_slide(C["white"])
    title(
        s,
        "A heterogeneous fleet, modelled honestly.",
        "Node classes differ because the work differs. The dispatcher matches a task to a node that"
        " can actually run it, then leases it.",
        False,
        15,
        tag="FLEET",
    )
    nodes = [
        (
            "macOS HOST",
            "launchd-managed host\ninstall; Apple toolchain\nwork stays native.",
            C["blue"],
        ),
        (
            "LINUX NODE",
            "native steward plus\ncontainerized execution\nunder Docker Engine.",
            C["orange"],
        ),
        (
            "KUBERNETES",
            "one orchestrator folds\nclaim, launch, and stuck-\nJob reconciliation.",
            C["green2"],
        ),
        (
            "HGX SESSION",
            "bounded elastic capacity,\nonboarded by explicit\noperator receipt.",
            C["green"],
        ),
    ]
    shape(s, "roundRect", 70, 296, 1140, 68, C["ink"], radius="rounded-xl")
    text(
        s,
        "DISPATCHER  ·  capability match, then lease",
        90,
        314,
        700,
        32,
        18,
        C["white"],
        True,
        fit=True,
    )
    for i, (head, body, color) in enumerate(nodes):
        x = 70 + i * 290
        down_arrow(s, x + 134, 364, 398, color)
        shape(
            s,
            "roundRect",
            x,
            398,
            272,
            200,
            C["fog"],
            radius="rounded-xl",
            stroke=color,
            strokeWidth=2,
        )
        text(s, head, x + 20, 420, 232, 30, 17, C["ink"], True, fit=True)
        line(s, x + 20, 456, 46, 4, color)
        text(s, body, x + 20, 476, 232, 100, 14, C["steel"], False, fit=True)
    caption_bar(
        s,
        "Capability, egress, and secrets are declared per node and per task - not assumed uniform.",
    )
    footer(s, 15)
    notes(
        s,
        [
            "macOS nodes are host installs under launchd by decision, not by omission; containerized execution is narrowed to Linux, where the container story is honest.",
            "Docker Engine / Moby is the single container runtime, so the sandbox policy has one enforcement path to reason about.",
            "HGX capacity is fungible and bounded: onboarding a provider session is an explicit operator action with a durable receipt, which keeps the machine-onboarding trust boundary intact.",
        ],
    )

    # 16 - multi-model routing
    s = add_slide(C["fog"])
    title(
        s,
        "Many models, one ordered route, one meter.",
        "Multi-model is a routing problem with a cost consequence. The ladder decides who does the"
        " work; the meter records what it cost.",
        False,
        16,
        tag="MODELS + MONEY",
    )
    text(s, "ROUTE LADDER  ·  ordered, capability-filtered", 70, 296, 600, 26, 14, C["steel"], True)
    ladder = ["opencode", "pi", "claude", "codex", "cursor"]
    for i, label in enumerate(ladder):
        x = 70 + i * 231
        fill = C["orange"] if i == 0 else C["white"]
        color = C["ink"]
        shape(
            s,
            "roundRect",
            x,
            330,
            213,
            62,
            fill,
            radius="rounded-lg",
            stroke=C["orange"] if i == 0 else C["line"],
            strokeWidth=2,
        )
        text(
            s, f"{i + 1}. {label}", x + 16, 348, 181, 26, 15, color, True, align="center", fit=True
        )
        if i < len(ladder) - 1:
            arrow(s, x + 213, 359, x + 231, C["steel"])
    shape(s, "roundRect", 70, 424, 550, 200, C["ink"], radius="rounded-xl")
    text(s, "RECORDED AT THE ROUTER", 96, 446, 400, 26, 15, C["orange2"], True)
    meter = [
        "one route event per model call",
        "input, output, and streaming state",
        "attributed to task, project, and agent",
    ]
    for i, item in enumerate(meter):
        dot(s, 96, 492 + i * 42, 10, C["orange"])
        text(s, item, 120, 484 + i * 42, 480, 34, 14, C["fog"], True, fit=True)
    shape(
        s,
        "roundRect",
        660,
        424,
        550,
        200,
        C["white"],
        radius="rounded-xl",
        stroke=C["blue"],
        strokeWidth=2,
    )
    text(s, "PRICED AT READ TIME", 686, 446, 400, 26, 15, C["blue"], True)
    priced = [
        "history re-prices without re-running",
        "cost per task, project, and outcome",
        "coverage gaps are measured, not averaged away",
    ]
    for i, item in enumerate(priced):
        dot(s, 686, 492 + i * 42, 10, C["blue"])
        text(s, item, 710, 484 + i * 42, 480, 34, 14, C["ink"], True, fit=True)
    footer(s, 16)
    notes(
        s,
        [
            "The ladder is ordered and filtered by capability and availability; the first entry is the current default coding route, and the order is a fleet contract rather than a per-worker preference.",
            "Recording happens where routing happens, which is the only place that sees every model call regardless of which agent made it. The router currently captures usage the caller reported; making the router itself the meter is a proposed decision, not a shipped one.",
            "Pricing at read time means a price-table change re-values history instead of invalidating it. Coverage is itself measured - the proposal that owns this quantified the unmetered fraction rather than averaging it away - so quote that figure from source-notes.md with its date.",
        ],
    )

    # 17 - evidence and measurement
    s = add_slide(C["ink"])
    title(
        s,
        "The fleet measures itself.",
        "Timing, spend, decisions, and provenance are recorded per step, which is what makes the"
        ' question "where did the day go?" answerable at all.',
        True,
        17,
        tag="EVIDENCE",
    )
    stage_flow(
        s,
        ["ROUTE EVENT", "ACTION EVENT", "GATE DECISION", "EVIDENCE POINTER", "AUDIT TRAIL"],
        300,
        dark=True,
    )
    cols = [
        ("PER TASK", "state history, owner,\nlease, dependencies"),
        ("PER STEP", "timing, tokens,\nexit decision"),
        ("PER RELEASE", "provenance pin and\nevidence closure"),
    ]
    for i, (head, body) in enumerate(cols):
        x = 70 + i * 386
        shape(
            s,
            "roundRect",
            x,
            396,
            368,
            160,
            C["panel"],
            radius="rounded-xl",
            stroke=C["blue"],
            strokeWidth=1,
        )
        text(s, head, x + 24, 418, 320, 28, 16, C["blue"], True, fit=True)
        text(s, body, x + 24, 452, 320, 76, 15, C["fog"], False, fit=True)
    text(
        s,
        "Evidence indexes what happened. It does not, by itself, authorize a release.",
        70,
        584,
        1140,
        40,
        20,
        C["orange2"],
        True,
        align="center",
        fit=True,
    )
    footer(s, 17, True)
    notes(
        s,
        [
            "The distinction on the bottom line matters and is easy to get wrong in a deck: the ledger and its pointers are an index of evidence, while acceptance and publication remain explicit decisions by an authorized actor.",
            "Action events arrive as a normalized stream from sandboxed execution, so 'what did the agent actually do' is not reconstructed from a transcript.",
            "Ledger census and token figures quoted in conversation must carry the date they were measured; source-notes.md holds the measured values and their dates.",
        ],
    )

    # 18 - scale
    s = add_slide(C["white"])
    title(
        s,
        "Scale of the implementation.",
        "Surface area, not maturity: what exists in the tree this package was audited against.",
        False,
        18,
        tag="SCALE",
    )
    facts = [
        ("221", "control-plane\nmodules"),
        ("435", "HTTP\nroutes"),
        ("458", "CLI leaf\ncommands"),
        ("12", "task\nstates"),
    ]
    for i, (num, label) in enumerate(facts):
        x = 70 + i * 290
        text(s, num, x, 316, 272, 76, 52, C["orange"] if i == 0 else C["ink"], True, align="center")
        line(s, x + 96, 404, 80, 4, C["line"])
        text(s, label, x, 424, 272, 60, 17, C["steel"], True, align="center", fit=True)
    pairs = [
        ("5", "coding-agent routes on the ladder"),
        ("4", "object groups: project, task, agent, admin"),
        ("2", "dispatch targets: fleet nodes and Kubernetes"),
    ]
    for i, (num, label) in enumerate(pairs):
        x = 70 + i * 386
        shape(s, "roundRect", x, 500, 368, 110, C["fog"], radius="rounded-xl")
        text(s, num, x + 24, 526, 60, 56, 32, C["orange"], True)
        text(s, label, x + 96, 526, 250, 62, 15, C["ink"], True, fit=True)
    caption_bar(
        s,
        "Every figure here is a count from the audited tree, recorded with its source in source-notes.md.",
    )
    footer(s, 18)
    notes(
        s,
        [
            "These counts describe surface area only. They are not a maturity claim and must be re-measured, not carried forward, when this package is regenerated.",
            "The CLI is the object model: project, task, agent, and admin. The HTTP API, the Python client, and the MCP server are clients of that same surface rather than parallel implementations.",
        ],
    )

    # 19 - stated honestly
    s = add_slide(C["fog"])
    title(
        s,
        "Stated honestly: implemented, decided, proposed.",
        "A deck that blurs these three is marketing. The distinction is kept here so an engineering"
        " reviewer can trust the rest of it.",
        False,
        19,
        tag="LIMITS",
    )
    columns = [
        (
            "IMPLEMENTED",
            C["green2"],
            [
                "Durable ledger, leases, recovery",
                "Sandboxed execution and egress policy",
                "Independent review gates",
                "Route events, priced at read time",
            ],
        ),
        (
            "DECIDED, NOT YET RUNTIME",
            C["orange"],
            [
                "Agent-initiated review scope",
                "Native steward plus containerized\nexecution on Linux",
            ],
        ),
        (
            "PROPOSED",
            C["steel"],
            [
                "Metering enforced at the router",
                "Route-search-path contract",
                "Task view as a graph",
                "Retrieval and extraction pipeline",
            ],
        ),
    ]
    for i, (head, color, items) in enumerate(columns):
        x = 70 + i * 386
        shape(
            s,
            "roundRect",
            x,
            300,
            368,
            320,
            C["white"],
            radius="rounded-xl",
            stroke=color,
            strokeWidth=2,
        )
        text(s, head, x + 24, 322, 320, 50, 16, C["ink"], True, fit=True)
        line(s, x + 24, 376, 46, 4, color)
        for j, item in enumerate(items):
            dot(s, x + 24, 408 + j * 56, 10, color)
            text(s, item, x + 48, 400 + j * 56, 296, 48, 14, C["steel"], True, fit=True)
    footer(s, 19)
    notes(
        s,
        [
            "Placement of each item traces to the architecture-decision record that owns it, with that record's own status line as the authority - see source-notes.md.",
            "'Decided, not yet runtime' is the category most decks omit. An accepted decision whose implementation is deferred is neither a shipped capability nor a proposal.",
        ],
    )

    # 20 - close
    s = add_slide(C["ink"])
    pill(s, BRAND, 70, 48, 148, C["white"], C["ink"])
    text(s, "Reuse the stack.\nOwn the truth.", 70, 150, 700, 160, 46, C["white"], True, fit=True)
    line(s, 74, 336, 120, 8, C["orange"])
    text(
        s,
        "NVIDIA OpenShell for confinement. NeMo Relay for observability. HGX for capacity. Postgres,"
        " Kubernetes, and the best available coding agents for execution. AgentFabric for the one"
        " thing none of them do: durable, auditable operational truth across the whole fleet.",
        74,
        366,
        660,
        170,
        19,
        C["fog"],
        False,
        fit=True,
    )
    steps = [
        ("START", "one repository,\none project, one\nreal workload"),
        ("GROW", "add node classes\nand coding routes\nas work demands"),
        ("MEASURE", "read cost and\nthroughput before\nadding capacity"),
    ]
    for i, (head, body) in enumerate(steps):
        y = 300 + i * 116
        shape(
            s,
            "roundRect",
            790,
            y,
            420,
            100,
            C["panel"],
            radius="rounded-xl",
            stroke=C["orange"],
            strokeWidth=1,
        )
        text(s, head, 814, y + 20, 120, 26, 15, C["orange"], True)
        text(s, body, 940, y + 18, 250, 66, 13, C["fog"], False, fit=True)
    text(s, str(TOTAL), 1162, 698, 48, 16, 12, C["muted"], True, align="right")
    notes(
        s,
        [
            "Close on the adoption path, not the vision: one project and one real workload first, because the fabric's value shows up when a second agent has to trust the first one's output.",
            "For mechanism questions, return to slides 7 through 10 (reuse and containment) and 12 through 17 (actors, lifecycle, coordination, fleet, routing, evidence).",
            "The authority for every claim in this deck is source-notes.md; the narrative member carries the same claims in prose at greater depth.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    return OUT


def _capability_manifest(pptx_path: Path) -> dict:
    return {
        "schema": "agentfabric/document-pair-manifest@1",
        "ecosystem": "google-workspace",
        "members": {
            "presentation": {
                "local_artifact": str(pptx_path),
                "published_location": None,
                "publication_authorized": False,
                "access": {
                    "audience": "organization",
                    "principals": [],
                    "permission": "view",
                    "link_sharing": "organization-restricted",
                },
            }
        },
        "authoring_package": {
            "root": str(SOURCE),
            "elements": {
                "deck_specification": "deck-specification.md",
                "narrative_specification": "narrative-specification.md",
                "factual_ledger": "source-notes.md",
                "generation_prompts": "prompts",
                "build_source": "build_deck.py",
                "regeneration_entry_point": "regenerate.sh",
                "deliverable_links": "current-deliverables.md",
                "qa_record": "qa-ledger.md",
            },
        },
    }


def main() -> None:
    out = build()
    slide_count = len(Presentation(str(out)).slides)
    print(f"built {slide_count} slides -> {out}")
    manifest_path = _obj_dir() / "agentfabric-overview" / "capability-manifest.json"
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(
        json.dumps(_capability_manifest(out), indent=2) + "\n", encoding="utf-8"
    )
    print(f"capability manifest -> {manifest_path}")


if __name__ == "__main__":
    main()
